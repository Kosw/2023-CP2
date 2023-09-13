from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()
title_slide_layout = prs. slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

title = slide.placeholders[0]
subtitle = slide.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "ComputerProgramming 2 First time."

tf = subtitle.text_frame
p = tf.add_paragraph()
p.text = "Add Paragraph"

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(title_slide_layout)

left = top = Inches(1)
width = height = Inches(5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame

tf.text = ("This is text inside a textbox")

p = tf.add_paragraph()
p.text = "This is second paragraph that's bold"
p.font.bold = True

p = tf.add_paragraph()
p.text = "This is third paragraph that's big"
p.font.size = Pt(20)

img_path = 'pixel_image.png'

slide = prs.slides.add_slide(blank_slide_layout)
left = top = Inches(1)
pic = slide.shapes.add_picture(img_path, left, top)

left = Inches(5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

prs. save("test.pptx")